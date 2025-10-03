import io
import os.path
import numpy as np
import pymupdf
import easyocr
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from sklearn.mixture import GaussianMixture
from sklearn.model_selection import GridSearchCV


class PDFtoPPTXConverter:

    def __init__(
        self,
        default_font: str | None = None,
        enable_ocr: bool = False,
        enforce_default_font: bool = False,
        lang: list[str] = None,
    ) -> None:
        """
        Args:
            default_font: Шрифт для OCR текста.
            enable_ocr: Использовать ли OCR. По умолчанию False.
            enforce_default_font: Использовать ли дефолтный шрифт всегда.
            lang: Список языков для EasyOCR (например ['ru', 'en']).
        """
        self.default_font = default_font
        self.enable_ocr = enable_ocr
        self.enforce_default_font = enforce_default_font
        
        if lang is None:
            lang = ['ru', 'en']
        
        self.ocr_reader = easyocr.Reader(lang, gpu=True) if self.enable_ocr else None

    def convert(self, input_file_path: str, output_file_path: str) -> None:
        """Конвертация PDF в PPTX."""
        
        self._validate_arguments(input_file_path, output_file_path)
        
        pptx_output = Presentation()
        
        with pymupdf.open(input_file_path) as pdf_document:
            # Определяем размер слайда из первой страницы
            first_page = pdf_document[0]
            slide_width = first_page.rect.width / 72.0
            slide_height = first_page.rect.height / 72.0
            
            pptx_output.slide_width = Inches(slide_width)
            pptx_output.slide_height = Inches(slide_height)
            
            all_contents = [
                self._get_page_contents(pdf_page) for pdf_page in pdf_document
            ]
            
            # Проверяем, отсканирован ли документ
            scanned_document = self.enable_ocr and all(
                (not page_contents["text_blocks"]) for page_contents in all_contents
            )
            
            self._construct_pptx(
                pdf_document, all_contents, pptx_output, scanned_document
            )
        
        output_directory = os.path.dirname(output_file_path)
        if output_directory and not os.path.exists(output_directory):
            os.makedirs(output_directory)
        
        pptx_output.save(output_file_path)

    def _get_page_contents(self, pdf_page):
        """Получение всего редактируемого контента со страницы."""
        page_contents = {}
        page_contents["text_blocks"] = self._get_page_text(pdf_page)
        page_contents["drawings"] = self._get_page_drawings(pdf_page)
        page_contents["images"] = self._get_page_images(pdf_page)
        page_contents["xref_smask_map"] = self._get_page_xref_smask_map(pdf_page)
        return page_contents

    def _get_page_text(self, pdf_page):
        """Получение всех текстовых блоков."""
        blocks = pdf_page.get_textpage().extractDICT()["blocks"]
        text_blocks = [block for block in blocks if block["type"] == 0]
        return text_blocks

    def _get_page_drawings(self, pdf_page):
        """Получение всех рисунков."""
        return pdf_page.get_drawings()

    def _get_page_images(self, pdf_page):
        """Получение всех изображений."""
        return pdf_page.get_image_info(xrefs=True)

    def _get_page_xref_smask_map(self, pdf_page):
        """Создание словаря smask для изображений."""
        return {item[0]: item[1] for item in pdf_page.get_images()}

    def _get_page_ocr_text(self, pdf_page):
        """OCR страницы с помощью EasyOCR."""
        page_pixmap = pdf_page.get_pixmap(dpi=300)
        page_bytes = page_pixmap.tobytes()
        
        # Конвертируем в PIL Image
        pil_image = Image.frombytes("RGB", (page_pixmap.width, page_pixmap.height), page_bytes)
        
        # EasyOCR работает с numpy array
        image_array = np.array(pil_image)
        ocr_result = self.ocr_reader.readtext(image_array)
        
        return ocr_result

    def _add_text_block_to_slide(self, text_block, slide):
        """Добавление текстового блока на слайд."""
        for line in text_block["lines"]:
            for span in line["spans"]:
                line_text = span["text"]
                line_font_size = span["size"]
                
                if line_font_size < 1:
                    continue
                    
                line_font_name = span["font"]
                line_font_color = span["color"]
                line_font_is_italic = bool(span["flags"] & 2**1)
                line_font_is_bold = bool(span["flags"] & 2**4)
                
                x0, y0, x1, y1 = span["bbox"]
                
                text_box_left = Inches(x0 / 72.0)
                text_box_top = Inches(y0 / 72.0)
                text_box_width = Inches((x1 - x0) / 72.0)
                text_box_height = Inches((y1 - y0) / 72.0)
                
                new_text_box = slide.shapes.add_textbox(
                    text_box_left, text_box_top, text_box_width, text_box_height
                )
                
                text_frame = new_text_box.text_frame
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                text_frame.margin_bottom = 0
                text_frame.margin_left = 0
                text_frame.margin_right = 0
                text_frame.margin_top = 0
                
                p = text_frame.paragraphs[0]
                p.text = line_text
                p.font.size = Pt(line_font_size)
                p.font.name = (
                    self.default_font
                    if self.enforce_default_font and self.default_font
                    else line_font_name
                )
                p.font.color.rgb = RGBColor.from_string(f"{line_font_color:06X}")
                p.font.italic = line_font_is_italic
                p.font.bold = line_font_is_bold
                p.alignment = PP_ALIGN.LEFT

    def _add_image_to_slide(self, pdf_doc, image, smask, bbox, slide):
        """Добавление изображения на слайд."""
        xref = image["xref"]
        base_image = pymupdf.Pixmap(pdf_doc, xref)
        
        if not base_image.colorspace:
            return
        
        if smask:
            smask_pixmap = pymupdf.Pixmap(pdf_doc, smask)
            if base_image.alpha:
                base_image.set_alpha(smask_pixmap.samples)
            else:
                base_image = pymupdf.Pixmap(base_image, smask_pixmap)
        
        transformation_matrix = pymupdf.Matrix(image["transform"])
        transformed_image_bytes = self._transform_image(
            base_image, bbox, transformation_matrix
        )
        
        left = Inches(bbox[0] / 72.0)
        top = Inches(bbox[1] / 72.0)
        width = Inches((bbox[2] - bbox[0]) / 72.0)
        
        slide.shapes.add_picture(
            io.BytesIO(transformed_image_bytes), left, top, width=width
        )

    def _add_drawing_to_slide(self, drawing, bbox, pdf_page_width, pdf_page_height, slide):
        """Добавление рисунка на слайд."""
        temp_pdf = pymupdf.open()
        temp_page = temp_pdf.new_page(width=pdf_page_width, height=pdf_page_height)
        
        shape = temp_page.new_shape()
        for item in drawing["items"]:
            if item[0] == "l":
                shape.draw_line(item[1], item[2])
            elif item[0] == "re":
                shape.draw_rect(item[1])
            elif item[0] == "qu":
                shape.draw_quad(item[1])
            elif item[0] == "c":
                shape.draw_bezier(item[1], item[2], item[3], item[4])
        
        kwargs_defaults = {
            "even_odd": True,
            "closePath": True,
            "lineJoin": 0,
            "width": 1,
            "stroke_opacity": 1,
            "fill_opacity": 1,
        }
        keys = ["fill", "color", "dashes", "even_odd", "closePath", "lineJoin", "width", "stroke_opacity", "fill_opacity"]
        kwargs = {
            key: (drawing.get(key) if drawing.get(key) is not None else kwargs_defaults.get(key))
            for key in keys
        }
        kwargs["lineCap"] = (
            max(drawing.get("lineCap")) if drawing.get("lineCap") is not None else 0
        )
        shape.finish(**kwargs)
        shape.commit()
        
        base_image = temp_page.get_pixmap(clip=bbox, dpi=300, alpha=True)
        try:
            image_bytes = base_image.tobytes()
            left = Inches(bbox[0] / 72.0)
            top = Inches(bbox[1] / 72.0)
            width = Inches((bbox[2] - bbox[0]) / 72.0)
            slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width=width)
        except:
            pass

    def _add_ocr_to_slide(self, ocr_line, fontsize, slide):
        """Добавление OCR текста на слайд."""
        # EasyOCR возвращает: (bbox, text, confidence)
        bbox, text, confidence = ocr_line
        
        # bbox это список из 4 точек [[x1,y1], [x2,y2], [x3,y3], [x4,y4]]
        scaling_factor = 0.24
        
        left = min(point[0] for point in bbox) * scaling_factor
        top = min(point[1] for point in bbox) * scaling_factor
        right = max(point[0] for point in bbox) * scaling_factor
        bottom = max(point[1] for point in bbox) * scaling_factor
        
        x, y = left, top
        width, height = right - left, bottom - top
        
        text_box = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(width), Pt(height))
        text_frame = text_box.text_frame
        text_frame.text = text
        
        text_frame.margin_bottom = 0
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(fontsize)
            if self.default_font:
                paragraph.font.name = self.default_font

    def _construct_pptx(self, pdf_document, all_contents, pptx_output, scanned_document):
        """Конструирование PPTX из PDF."""
        if scanned_document:
            ocr_results = [
                self._get_page_ocr_text(pdf_page) for pdf_page in pdf_document
            ]
            detected_sizes = [
                self._get_ocr_detected_font_size(line)
                for ocr_result in ocr_results
                for line in ocr_result
            ]
            suitable_sizes = self._get_suitable_sizes(detected_sizes)
        
        for page_num, pdf_page in enumerate(pdf_document):
            pdf_page_width = pdf_page.rect.width
            pdf_page_height = pdf_page.rect.height
            
            slide_layout = pptx_output.slide_layouts[6]
            slide = pptx_output.slides.add_slide(slide_layout)
            
            page_content = all_contents[page_num]
            drawings = page_content["drawings"]
            drawing_bboxes = [drawing["rect"] for drawing in drawings]
            images = page_content["images"]
            image_bboxes = [pymupdf.Rect(image["bbox"]) for image in images]
            xref_smask_map = page_content["xref_smask_map"]
            all_shapes = drawings + images
            all_bboxes = drawing_bboxes + image_bboxes
            
            indices = self._sort_shapes_by_area(all_bboxes)
            
            for index in indices:
                shape = all_shapes[index]
                bbox = all_bboxes[index]
                
                if index < len(drawings):
                    if scanned_document and self._drawing_is_character(shape):
                        continue
                    self._add_drawing_to_slide(
                        shape, bbox, pdf_page_width, pdf_page_height, slide
                    )
                else:
                    xref = shape["xref"]
                    if xref == 0:
                        continue
                    smask = xref_smask_map[xref]
                    self._add_image_to_slide(pdf_document, shape, smask, bbox, slide)
            
            text_blocks = page_content["text_blocks"]
            for text_block in text_blocks:
                self._add_text_block_to_slide(text_block, slide)
            
            if scanned_document:
                page_ocr_results = ocr_results[page_num]
                for line_num, ocr_line in enumerate(page_ocr_results):
                    fontsize = suitable_sizes[
                        sum(len(page_ocr) for page_ocr in ocr_results[:page_num]) + line_num
                    ]
                    self._add_ocr_to_slide(ocr_line, fontsize, slide)

    @staticmethod
    def _sort_shapes_by_area(all_bboxes):
        """Сортировка фигур по площади."""
        indices = sorted(
            range(len(all_bboxes)),
            key=lambda i: all_bboxes[i].get_area(),
            reverse=True,
        )
        return indices

    @staticmethod
    def _get_ocr_detected_font_size(ocr_line):
        """Определение размера шрифта из OCR результата."""
        scaling_factor = 0.24
        
        # EasyOCR: bbox это список из 4 точек
        bbox = ocr_line[0]
        left = min(point[0] for point in bbox) * scaling_factor
        top = min(point[1] for point in bbox) * scaling_factor
        right = max(point[0] for point in bbox) * scaling_factor
        bottom = max(point[1] for point in bbox) * scaling_factor
        
        height = bottom - top
        width = right - left
        
        fontsize = np.clip(min(height, width), 1, 64)
        return fontsize

    @staticmethod
    def _gmm_bic_score(estimator, X):
        """BIC score для GridSearchCV."""
        return -estimator.bic(X)

    @staticmethod
    def _get_suitable_sizes(detected_sizes):
        """Определение подходящих размеров шрифта с помощью GMM."""
        if len(detected_sizes) < 5:
            return detected_sizes
        
        X = np.asarray(detected_sizes).reshape((-1, 1))
        
        param_grid = {"n_components": range(1, 7)}
        grid_search = GridSearchCV(
            GaussianMixture(init_params="k-means++"),
            param_grid=param_grid,
            scoring=PDFtoPPTXConverter._gmm_bic_score,
        )
        grid_search.fit(X)
        
        best_estimator = grid_search.best_estimator_
        labels = best_estimator.predict(X)
        suitable_sizes = list(map(lambda i: best_estimator.means_[i], labels))
        return suitable_sizes

    @staticmethod
    def _transform_image(base_image, bbox, transformation_matrix):
        """Трансформация изображения."""
        base_image = pymupdf.Pixmap(pymupdf.Colorspace(pymupdf.CS_RGB), base_image)
        
        if not base_image.alpha:
            alpha = bytearray(255 for _ in range(base_image.w * base_image.h))
            base_image = pymupdf.Pixmap(base_image, 1)
            base_image.set_alpha(alpha)
        
        image_data = base_image.samples
        pil_image = Image.frombytes("RGBA", (base_image.w, base_image.h), image_data)
        
        scaling_factor = 4
        
        shrink = pymupdf.Matrix(
            scaling_factor / base_image.w, 0, 0, scaling_factor / base_image.h, 0, 0
        )
        translation = pymupdf.Matrix(1, 0, 0, 1, -bbox[0], -bbox[1])
        transformation_matrix = shrink * transformation_matrix * translation
        transformation_matrix.invert()
        
        pil_image = pil_image.transform(
            (
                scaling_factor * int(bbox[2] - bbox[0]),
                scaling_factor * int(bbox[3] - bbox[1]),
            ),
            Image.AFFINE,
            (
                transformation_matrix.a,
                transformation_matrix.c,
                transformation_matrix.e,
                transformation_matrix.b,
                transformation_matrix.d,
                transformation_matrix.f,
            ),
            resample=Image.Resampling.BICUBIC,
        )
        
        output_bytes = io.BytesIO()
        pil_image.save(output_bytes, format="PNG")
        return output_bytes.getvalue()

    @staticmethod
    def _drawing_is_character(drawing):
        """Определение, является ли рисунок символом."""
        command_length_threshold = 6
        return all([
            not drawing["closePath"],
            drawing["type"] == "f",
            len(drawing["items"]) > command_length_threshold,
        ])

    @staticmethod
    def _validate_arguments(input_file_path, output_file_path):
        """Валидация аргументов."""
        if not os.path.isfile(input_file_path):
            raise FileNotFoundError(f"Файл {input_file_path} не существует.")
        
        input_file_extension = os.path.splitext(input_file_path)[1]
        if input_file_extension != ".pdf":
            raise ValueError(
                f"Неверное расширение входного файла. Ожидается .pdf, получено {input_file_extension}"
            )
        
        output_file_extension = os.path.splitext(output_file_path)[1]
        if output_file_extension != ".pptx":
            raise ValueError(
                f"Неверное расширение выходного файла. Ожидается .pptx, получено {output_file_extension}"
            )


# Главная программа
if __name__ == "__main__":
    # Запрашиваем имя входного файла
    input_file = input("Введите имя входного файла (PDF): ")
    
    # Запрашиваем имя выходного файла
    output_file = input("Введите имя выходного файла (PPTX): ")
    
    # Создаем конвертер
    converter = PDFtoPPTXConverter(
        enable_ocr=True,
        lang=['ru', 'en']
    )
    
    # Выполняем конвертацию
    print(f"Начинаем конвертацию {input_file} в {output_file}...")
    converter.convert(input_file, output_file)
    print("Конвертация завершена!")
